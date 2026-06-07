#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
全駒煙（協力詰・40駒）フェアリー全国大会 特別講演スライド生成。
スピーカーノート付き。出力: zenkoma_talk.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette ----
INK   = RGBColor(0x1A, 0x1A, 0x1A)
SUB   = RGBColor(0x55, 0x55, 0x55)
ACC   = RGBColor(0xB0, 0x3A, 0x2E)   # 落ち着いた朱
ACC2  = RGBColor(0x2E, 0x5A, 0x88)   # 藍
LIGHT = RGBColor(0xEC, 0xE6, 0xDC)   # 和紙系
BAND  = RGBColor(0xC8, 0x9B, 0x3C)   # 有望帯の金
SHADE = RGBColor(0xD8, 0xD2, 0xC8)

FONT = "Meiryo"   # 日本語が無難に出るフォント（無ければ閲覧側で代替される）

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def txt(slide, l, t, w, h, text, size=18, bold=False, color=INK,
        align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def bullets(slide, l, t, w, h, items, size=18, color=INK, gap=6, lead=""):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (lvl, s, *opt) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        mark = "● " if lvl == 0 else ("– " if lvl == 1 else "· ")
        r = p.add_run()
        r.text = mark + s
        r.font.size = Pt(size - lvl * 2)
        r.font.name = FONT
        r.font.color.rgb = opt[0] if opt else color
        r.font.bold = (lvl == 0 and len(opt) > 1)
    return tb


def bar(slide, color=ACC):
    """左の細い縦アクセント"""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                Inches(0.18), SH)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    return sp


def header(slide, kicker, title, color=ACC):
    bar(slide, color)
    txt(slide, Inches(0.55), Inches(0.42), Inches(12), Inches(0.4),
        kicker, size=14, bold=True, color=color)
    txt(slide, Inches(0.5), Inches(0.78), Inches(12.3), Inches(1.0),
        title, size=30, bold=True, color=INK)


def placeholder_diagram(slide, l, t, w, h, label):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = LIGHT
    sp.line.color.rgb = SUB; sp.line.width = Pt(0.75)
    tf = sp.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(14); r.font.color.rgb = SUB; r.font.name = FONT
    return sp


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


# =====================================================================
# 1. タイトル
# =====================================================================
s = add_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x22)
bg.line.fill.background()
txt(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.4),
    "全駒煙 ― 協力詰、極限への到達", size=44, bold=True,
    color=RGBColor(0xF4, 0xF0, 0xE8))
txt(s, Inches(0.95), Inches(3.45), Inches(11), Inches(0.8),
    "盤上40枚すべてが消えてゆく、ひとつの煙詰", size=22,
    color=BAND)
txt(s, Inches(0.95), Inches(5.6), Inches(11), Inches(1.0),
    "荻絵香木（ogiekako）\nフェアリー全国大会 前夜祭 ／ 特別講演", size=18,
    color=RGBColor(0xC8, 0xC2, 0xB8))
note(s, """
つかみ。最初に一号局を盤で実際に並べる（または投影して解いて見せる）。
「将棋の駒を一組すべて盤に載せた煙詰です。協力詰では、これが世界で最初の一局＝一号局です。」
ここで40枚が煙のように消え、最後は玉＋1枚で詰む様子を1手ずつ見せると一番つかめる。
講演全体は約30〜40分。技術自慢ではなく『手では届かない場所を見にいった話』として語る。
""")

# =====================================================================
# 2. 自己紹介と立ち位置
# =====================================================================
s = add_slide()
header(s, "1 / 自己紹介", "私はもともと、手で協力詰を作ってきた人間です", ACC2)
bullets(s, Inches(0.6), Inches(1.95), Inches(6.4), Inches(4.6), [
    (0, "長年、協力詰（helpmate）を手作業で創作", INK, True),
    (0, "単玉協力詰の最長手数を手で大幅更新", INK, True),
    (1, "「ミトコンドリア」3933手 ― 一手ずつ手で配置を探した"),
    (1, "詰将棋パラダイス 2026年6月号「神無一族の氾濫」で発表"),
    (0, "代表作のひとつ「チェイン」"),
    (0, "── その私が、手では決して届かない場所を", ACC, True),
    (0, "　　 見にいくために作った道具の話を、きょうはします", ACC, True),
])
placeholder_diagram(s, Inches(7.3), Inches(2.0), Inches(5.4), Inches(4.5),
    "［局面図を貼る］\n単玉最長 3933手「ミトコンドリア」\n\n"
    "G+p2GSGSG/1p2S1S1P/1l1L1N1N+R/\nk1+P3PP1/+r1b3N2/3b4N/\n1P6+p/2P6/3PPLL2 b P6p 1")
note(s, """
ここが講演で最も大事な前置き。AI創作への懸念に一言も触れずに、構図そのものを無効化する。
語り口の例：
「私はもともと、手で協力詰を作ってきた人間です。たとえばこれは単玉の最長手数を更新した
『ミトコンドリア』、3933手。一手ずつ手で配置を探した作品です（詰パラ6月号で発表）。
── きょうお話しするのは、その私が、手では決して届かない場所を見にいくために作った道具の話です。」

ポイント：
- 自分が分野の内部の作家であることを、事実として淡々と置くだけでよい。防御も自己弁護もしない。
- 「機械で分野を侵略した」という見方は、この実作実績の前では成立しない。論じる必要すらない。
- ミトコンドリア1本に絞るか、チェインも併用するかは好み。1本に絞ると入りがシャープ。
""")

# =====================================================================
# 3. 全駒煙とは / なぜ40が究極か
# =====================================================================
s = add_slide()
header(s, "2 / 全駒煙とは", "なぜ「40枚」が究極なのか")
bullets(s, Inches(0.6), Inches(1.95), Inches(7.0), Inches(4.6), [
    (0, "煙詰：多数の駒が「煙のように」消え、玉＋1枚＝2枚で詰む", INK, True),
    (0, "全駒煙：初形が将棋一組すべて＝盤上40枚・持駒ゼロ", INK, True),
    (1, "王2 飛2 角2 金4 銀4 桂4 香4 歩18 ＝ 計40"),
    (1, "41枚はあり得ない ＝ 理論上の最大"),
    (0, "だからこれは「佳作の量産」ではない", ACC, True),
    (1, "唯一の極限を『在るのか／一意か』と問う話", ACC),
    (1, "性格としては作曲よりも“発見”に近い", ACC),
])
placeholder_diagram(s, Inches(7.9), Inches(2.0), Inches(4.8), Inches(4.5),
    "［一号局の局面図を貼る］\n\n40枚 → … → 玉＋1枚\n「煙のように消える」")
note(s, """
全駒煙の定義と、40が天井であることを明確に。
- 40 = 一組全部・持駒ゼロ。これより多い初形は物理的に存在しない。
- ここで『極限・記録の探索』という性格づけを与えておく。これが後段でAI懸念をかわす伏線になる
  （恣意的に作品を量産したのではなく、唯一の極限が在るかを確かめた、という枠組み）。
- 残る問いは「40を何手で、何通り作れるか」だけになる、と次章以降へ橋渡し。
""")

# =====================================================================
# 4. ルールの垣根
# =====================================================================
s = add_slide()
header(s, "3 / ルールの垣根", "協力詰は、普通の詰将棋とは別の問題", ACC2)
bullets(s, Inches(0.6), Inches(1.95), Inches(12.2), Inches(4.6), [
    (0, "今大会では、正調煙詰の自動生成（広瀬稔・馬屋原剛 両氏）も登場", INK, True),
    (0, "本作は協力詰（helpmate）― 攻方・受方の双方が協力して詰みに向かう", INK, True),
    (1, "探索木に AND/OR（攻防の min/max）構造が無い"),
    (1, "→ 証明数探索（df-pn）は原理的に効かない ＝ 別系統の技術が要る"),
    (0, "ルールが違えば、機械の作り方も全く違う", ACC, True),
    (1, "正調煙詰と協力詰、二つの道から同じ「煙」へ近づいている", ACC),
])
note(s, """
もう一方の発表（広瀬・馬屋原両氏の正調煙詰自動生成）と自然に並置し、共同出題の伏線を張る。
技術的な核心を一枚で：
- 協力詰は双方が協力するので、通常詰将棋の『攻める側 vs 逃げる側』という対立（min/max）が無い。
- そのため、通常詰将棋ソルバの主力である証明数探索（df-pn）がそのままでは使えない。
- だから「別の技術」が必要になる、という必然性をここで示す（次章の動機づけ）。
深入りは禁物。この一枚で言い切って次へ。
責任者メッセージ『ルールの垣根を越え、変革期にいる』にここで触れてもよい。
""")

# =====================================================================
# 5. なぜ届かないか（探索誘導の問題）
# =====================================================================
s = add_slide()
header(s, "4 / なぜ難しいか", "力任せでは、手でも計算機でも届かない")
bullets(s, Inches(0.6), Inches(1.95), Inches(6.3), Inches(4.6), [
    (0, "40枚に至る局面の数は 10^10〜10^11 規模", INK, True),
    (1, "厳密な全探索は、数百GB〜数TBでも途中で破綻"),
    (0, "これは「速度」の問題ではない", ACC, True),
    (0, "「どの局面を残すか＝探索の勘どころ」の問題", ACC, True),
    (1, "全部は見られない。見るべき所を選ぶ必要がある"),
    (1, "＝ 作家が盤面に対して持つ“嗅覚”を、機械に移すこと"),
])
# 右：指数爆発＋薄い有望帯の図
ox, oy, ow, oh = Inches(7.2), Inches(2.0), Inches(5.6), Inches(4.4)
fr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, ox, oy, ow, oh)
fr.fill.solid(); fr.fill.fore_color.rgb = RGBColor(0xFB, 0xF9, 0xF4)
fr.line.color.rgb = SUB; fr.line.width = Pt(0.75)
# 爆発する三角（影）
tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                         ox + Inches(0.5), oy + Inches(0.5),
                         Inches(4.6), Inches(3.0))
tri.rotation = 0
tri.fill.solid(); tri.fill.fore_color.rgb = SHADE
tri.line.fill.background()
txt(s, ox + Inches(2.0), oy + Inches(0.7), Inches(3.2), Inches(0.6),
    "全探索＝指数爆発\n（到達不能）", size=12, color=SUB, align=PP_ALIGN.CENTER)
# 薄い有望帯（金の細帯）
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          ox + Inches(0.5), oy + Inches(3.05),
                          Inches(4.6), Inches(0.22))
band.fill.solid(); band.fill.fore_color.rgb = BAND
band.line.fill.background()
txt(s, ox + Inches(0.5), oy + Inches(3.32), Inches(4.8), Inches(0.7),
    "← beam が残す薄い“有望帯”\n（作家の嗅覚＝特徴量で選別）",
    size=12, color=ACC, align=PP_ALIGN.LEFT)
txt(s, ox + Inches(0.3), oy + Inches(3.95), Inches(5.0), Inches(0.4),
    "深さ（駒が増える方向）→", size=11, color=SUB)
note(s, """
ここが『機械が力ずくで塗りつぶした』という像を崩す最重要スライド。
- 全駒煙の探索空間は天文学的（10^10〜11）。厳密全幅は大規模マシンでも破綻する。
- 中核メッセージ：『これは速度の問題ではなく、どの局面を残すか＝探索の勘どころの問題』。
- だから、力任せでは解けない。人間（作家）の“行けそう／行けなさそう”という嗅覚を機械に
  移植して初めて届いた。── この物語の起点をここで宣言する。
右の図：上に向かって指数的に広がる空間（影）と、その中の薄い金色の帯（beamが残す有望帯）。
""")

# =====================================================================
# 6. 技術ハイライト① 構造の発見
# =====================================================================
s = add_slide()
header(s, "5 / 技術ハイライト①", "解は「長い背骨＋浅い冠」という形をしている", ACC2)
# 図：冠（扇）→ 合流点（箱）→ 背骨（縦線）→ 詰み（円）
cx = Inches(3.0)
# crown fan lines
import math
fan_top_y = Inches(1.9)
gate_y = Inches(3.3)
for dx in (-1.7, -1.0, -0.35, 0.35, 1.0, 1.7):
    ln = s.shapes.add_connector(2, Emu(int(cx) + Inches(dx)), fan_top_y, cx, gate_y)
    ln.line.color.rgb = ACC2; ln.line.width = Pt(1.5)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Emu(int(cx) + int(Inches(dx)) - int(Inches(0.1))),
                             Emu(int(fan_top_y) - int(Inches(0.1))),
                             Inches(0.2), Inches(0.2))
    dot.fill.solid(); dot.fill.fore_color.rgb = ACC2; dot.line.fill.background()
txt(s, Inches(0.4), Inches(1.45), Inches(5.2), Inches(0.4),
    "初形（40枚）← 2618局がここで枝分かれ＝浅い「冠」", size=12,
    color=ACC2, align=PP_ALIGN.CENTER)
# gateway box
gate = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Emu(int(cx) - int(Inches(1.3))), gate_y,
                          Inches(2.6), Inches(0.55))
gate.fill.solid(); gate.fill.fore_color.rgb = ACC; gate.line.fill.background()
tf = gate.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "合流点（gateway）"; r.font.size = Pt(13)
r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); r.font.name = FONT
# spine
spine = s.shapes.add_connector(2, cx, Emu(int(gate_y) + int(Inches(0.55))),
                               cx, Inches(5.9))
spine.line.color.rgb = INK; spine.line.width = Pt(3)
txt(s, Emu(int(cx) + int(Inches(0.25))), Inches(4.25), Inches(2.6), Inches(1.0),
    "1本の forced spine（背骨）\n全解が完全に一致する\n長い一本道", size=12, color=INK)
mate = s.shapes.add_shape(MSO_SHAPE.OVAL,
                          Emu(int(cx) - int(Inches(0.45))), Inches(5.9),
                          Inches(0.9), Inches(0.45))
mate.fill.solid(); mate.fill.fore_color.rgb = INK; mate.line.fill.background()
tf = mate.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "詰み"; r.font.size = Pt(11)
r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); r.font.name = FONT
bullets(s, Inches(6.5), Inches(2.1), Inches(6.3), Inches(4.3), [
    (0, "全2618解が必ず通る唯一の関門＝合流点", INK, True),
    (0, "多様性は、長い必然の一本道の上に乗った", INK),
    (0, "　浅い“冠”にすぎない", INK),
    (0, "→ 「どこを探せばいいか」は闇雲ではなく", ACC, True),
    (0, "　 この構造が教えてくれる", ACC, True),
])
note(s, """
技術ハイライトの第1の山。図解スライド。
- 全2618解を mate（詰み）から逆にたどって距離ごとに数えると、ほとんどの距離で局面は『1個』しかない。
  ＝ 全解が完全に一致する長い一本道（背骨／forced spine）。
- 根（初形）に近い側でだけ枝分かれする（冠／crown）。例：根で2618に開くが、その手前は数百〜数個。
- 全解が必ず通る最深の1点＝合流点（gateway）。ここを見つけたことが、探索を集中させる鍵になった。
聴衆へのメッセージ：『40枚の多様性は、必然の一本道に乗った浅い冠にすぎない。だから探索は闇雲ではなく、
構造に導かれて行える』。数値（d104→136, 根→2618 等）は口頭で軽く添える程度でよい。
""")

# =====================================================================
# 7. 技術ハイライト② 作家の直感の数値化
# =====================================================================
s = add_slide()
header(s, "5 / 技術ハイライト②", "作家の“直感”を、機械に教える")
bullets(s, Inches(0.6), Inches(1.95), Inches(7.2), Inches(4.6), [
    (0, "「行けそう／行けなさそう」という作家の感覚を", INK, True),
    (0, "　約85個の数値（特徴量）に翻訳", INK, True),
    (1, "玉の自由度・逃げ道、駒の散り具合、成駒の量 …"),
    (0, "機械はそれを手がかりに、薄い“有望帯”だけを", INK, True),
    (0, "　残しながら深部へ潜る（ビームサーチ）", INK, True),
    (0, "機械は人間の嗅覚の“代替”ではなく“延長”", ACC, True),
    (1, "力任せの計算ではなく、移植された勘で進む", ACC),
])
placeholder_diagram(s, Inches(8.1), Inches(2.0), Inches(4.6), Inches(4.4),
    "［模式図：盤面 → 約85の特徴量 → スコア］\n\n"
    "玉の自由度／逃走深さ\n配置の広がり（dispersion）\n成駒の量・位置 …\n\n"
    "↓\n\n有望な局面だけを残す")
note(s, """
技術ハイライトの第2の山。4章の物語（嗅覚の移植）をここで回収する。
- 作家が盤を見て感じる『この形は行けそう／だめそう』を、約85個の数値（特徴量）にした。
  例：玉の自由度・安全な逃げ場の数・逃走の深さ、駒の散らばり方、成駒の量や位置 など。
- 機械学習モデルがそれらから各局面のスコアを出し、ビームサーチが上位の薄い帯だけを残して深部へ進む。
- ここで明言：『機械は人間の嗅覚の代替ではなく延長である』。AI懸念への返答を、技術の語り口に溶かす。
専門的になりすぎないよう、特徴量は2〜3例だけ口頭で挙げる。モデルの内部（GBDT等）は名前だけでよい。
""")

# =====================================================================
# 8. 技術ハイライト③ 厳密性
# =====================================================================
s = add_slide()
header(s, "5 / 技術ハイライト③", "見つけた40枚は、すべて厳密に検証されている", ACC2)
bullets(s, Inches(0.6), Inches(1.95), Inches(12.2), Inches(4.6), [
    (0, "ビームは「どこを探すか」を絞るだけ", INK, True),
    (0, "見つけた各局面の一意性は、一つ残らず厳密に確認", INK, True),
    (1, "解がちょうど一本道であることをソルバで検証"),
    (1, "「それっぽいもの」ではなく、本物のユニーク協力詰"),
    (0, "計算規模：数百コア・数TB級のマシンで数時間（クラウド）", SUB),
    (0, "── 探索は近似でも、作品の正しさは厳密", ACC, True),
])
note(s, """
作家・解答者の信頼に直結するので必ず置く。
- ビームサーチは『探索範囲を絞る近似』だが、見つけた一つ一つの局面は、解が一意（協力手順が
  ちょうど一本）であることをソルバで厳密に検証している。だから本物。
- 留保も正直に：113手・2618局という“数”は近似探索が到達した範囲であり、『これが全ての40枚』
  『113手が真の最長』とまでは厳密証明していない。論文や発表で網羅性を主張しない。
  （聴衆が研究者寄りなら一言添えると誠実。一般寄りなら口頭で軽く。）
- 計算規模は1行で。自慢にしない。
""")

# =====================================================================
# 9. 発見されたもの ＋ 人の手
# =====================================================================
s = add_slide()
header(s, "6 / 見つかったもの", "そして、人の手が残る場所")
bullets(s, Inches(0.6), Inches(1.95), Inches(7.0), Inches(4.6), [
    (0, "113手・2618局 ― すべて一意・全駒盤上", INK, True),
    (0, "発表用の「美しい一局」を選ぶ", INK, True),
    (1, "成駒の少なさ、玉への密集度などで序列化"),
    (1, "上位を選び ― 最後に角を1マス寄せて手で仕上げた"),
    (0, "機械が候補を出し、選び・整えるのは人間", ACC, True),
    (1, "何を美しいと感じるかは、まだ人間の仕事", ACC),
])
placeholder_diagram(s, Inches(7.9), Inches(2.0), Inches(4.8), Inches(4.5),
    "［発表用の最終形・局面図を貼る］\n\n"
    "8G/6K1p/4S3+S/3Gp+PP+p+n/\n2B+pPP1p1/3Bk+p+p1P/\n1RRs+lgLN+P/2+PPNlNPG/\n+P+pL4S1 b - 1\n\n113手・一意")
note(s, """
ここは『山にはしない』が、一度だけ静かに置く重要な一枚。AI創作論への返答を独立章にせず、ここに溶かす。
- 機械は113手・2618局を出した。すべて一意で全駒盤上。
- だが『どれを見せるか』は機械には決められない。成駒の少なさ・玉への密集度などの“美の物差し”は
  人間が与えた。上位を選び、最後は作家が角を1マス寄せて手で仕上げ、より整った最終形にした。
- メッセージ：『機械が候補を出し、選び・整えるのは人間。何を美しいと感じるかは、まだ人間の仕事』。
  ここで一度だけ、静かに言って次へ。力まない。
""")

# =====================================================================
# 10. これから
# =====================================================================
s = add_slide()
header(s, "7 / これから", "ルールを越えて", ACC2)
bullets(s, Inches(0.6), Inches(2.1), Inches(12.2), Inches(4.2), [
    (0, "正調煙詰（自動生成）と協力詰、二つの道から同じ「煙」へ", INK, True),
    (0, "ルールの垣根を越えた共同出題", INK, True),
    (0, "私たちは道具を得て、手では届かなかった場所に手が届いた", ACC, True),
    (0, "そこで何を美しいと感じ、どれを選ぶかは ── まだ人間の仕事", ACC, True),
])
note(s, """
開いた終わり方。変革期の前夜祭にふさわしく。
- 二つの発表（正調煙詰の自動生成／協力詰の全駒煙）が、ルールを越えて出会った、という構図を示す。
- 共同出題が確定なら、ここをクライマックスにできる（具体的な出題内容に触れる）。
- 最後の一言は、得たもの（届かなかった場所への到達）と、人間に残るもの（美の判断）を並べて締める。
  断定ではなく『そうあってほしい』という願いの温度で。
""")

# =====================================================================
# 11. 結び
# =====================================================================
s = add_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x22)
bg.line.fill.background()
txt(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.4),
    "煙の果てに", size=40, bold=True, color=RGBColor(0xF4, 0xF0, 0xE8))
txt(s, Inches(0.95), Inches(4.0), Inches(11), Inches(1.2),
    "手では届かない場所を見にいくための、ひとつの道具の話でした。\nご清聴ありがとうございました。",
    size=20, color=BAND)
note(s, """
結び。質疑への入り口。
想定問答（準備しておくと安心）：
- 「全部の40枚を見つけたのか？」→ いいえ。近似探索が到達した範囲。網羅は今後の課題（厳密計算が要る）。
- 「これは創作か発見か？」→ 極限（40＝天井）の発見に近い。美の選択は人間が担っている。
- 「手作りの作家として、機械を使うことに葛藤は？」→ 正直な自分の言葉で。私は手で作り続けてもいる、と。
""")

out = "/home/ogiekako/fmrs/rust/zenkoma_talk.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
